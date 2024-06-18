import streamlit as st
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from fpdf import FPDF
from PIL import Image
from docx import Document
import io
import pypandoc
import pytesseract
import pandas as pd
from pptx import Presentation
from pdf2image import convert_from_path
import fitz  # This is the import for PyMuPDF
import tempfile
import os

# Inject the Google site verification meta tag into the head section
st.markdown(
    """
    <meta name="google-site-verification" content="fFXSjUHkzBdz0sWBJy2uAIHgwTTswufOA_x-qkSSqr8" />
    """,
    unsafe_allow_html=True
)

# Download Pandoc if it's not installed
try:
    pypandoc.get_pandoc_path()
except OSError:
    pypandoc.download_pandoc()

# Helper functions
def merge_pdfs(uploaded_files):
    merger = PdfMerger()
    for uploaded_file in uploaded_files:
        merger.append(uploaded_file)
    output = io.BytesIO()
    merger.write(output)
    merger.close()
    output.seek(0)
    return output

def split_pdf(uploaded_file, start_page, end_page):
    pdf_reader = PdfReader(uploaded_file)
    pdf_writer = PdfWriter()
    for page_num in range(start_page - 1, end_page):
        pdf_writer.add_page(pdf_reader.pages[page_num])
    output = io.BytesIO()
    pdf_writer.write(output)
    output.seek(0)
    return output

def compress_pdf(uploaded_file):
    pdf_reader = PdfReader(uploaded_file)
    pdf_writer = PdfWriter()
    for page in pdf_reader.pages:
        pdf_writer.add_page(page)
    output = io.BytesIO()
    pdf_writer.write(output)
    output.seek(0)
    return output

def images_to_pdf(uploaded_files):
    pdf = FPDF()
    for uploaded_file in uploaded_files:
        image = Image.open(uploaded_file)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_file:
            image.save(temp_file.name, 'JPEG')
            pdf.add_page()
            width, height = image.size
            aspect_ratio = width / height
            if width > height:
                pdf_width = 210
                pdf_height = pdf_width / aspect_ratio
            else:
                pdf_height = 297
                pdf_width = pdf_height * aspect_ratio
            x = (210 - pdf_width) / 2
            y = (297 - pdf_height) / 2
            pdf.image(temp_file.name, x, y, pdf_width, pdf_height)
    output = io.BytesIO()
    pdf_output = pdf.output(dest='S').encode('latin1')
    output.write(pdf_output)
    output.seek(0)
    return output

def word_to_pdf(uploaded_file):
    document = Document(uploaded_file)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for para in document.paragraphs:
        pdf.multi_cell(0, 10, para.text)
    output = io.BytesIO()
    pdf_output = pdf.output(dest='S').encode('latin1')
    output.write(pdf_output)
    output.seek(0)
    return output

def pdf_to_word(uploaded_file):
    pdf_reader = PdfReader(uploaded_file)
    doc = Document()
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text = page.extract_text()
        doc.add_paragraph(text)
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def unlock_pdf(uploaded_file, password):
    pdf_reader = PdfReader(uploaded_file)
    if pdf_reader.is_encrypted:
        pdf_reader.decrypt(password)
    pdf_writer = PdfWriter()
    for page in pdf_reader.pages:
        pdf_writer.add_page(page)
    output = io.BytesIO()
    pdf_writer.write(output)
    output.seek(0)
    return output

def protect_pdf(uploaded_file, password):
    pdf_reader = PdfReader(uploaded_file)
    pdf_writer = PdfWriter()
    for page in pdf_reader.pages:
        pdf_writer.add_page(page)
    pdf_writer.encrypt(password)
    output = io.BytesIO()
    pdf_writer.write(output)
    output.seek(0)
    return output

def ppt_to_pdf(uploaded_file):
    prs = Presentation(uploaded_file)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    for slide in prs.slides:
        pdf.add_page()
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = ""
                    for run in paragraph.runs:
                        text += run.text + "\n"
                    pdf.multi_cell(0, 10, text.encode('latin1', 'replace').decode('latin1'))

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        pdf.output(temp_file.name)
        temp_file_path = temp_file.name

    with open(temp_file_path, "rb") as f:
        pdf_bytes = f.read()

    output = io.BytesIO(pdf_bytes)
    output.seek(0)
    return output

def pdf_to_ppt(uploaded_file):
    pdf_reader = PdfReader(uploaded_file)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[5]
    for page_num in range(len(pdf_reader.pages)):
        slide = prs.slides.add_slide(blank_slide_layout)
        page = pdf_reader.pages[page_num]
        text = page.extract_text()
        txBox = slide.shapes.add_textbox(20, 20, 620, 440)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = text
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# Streamlit UI
st.set_page_config(page_title="PDF Tools", page_icon="📄", layout="wide")

# Add a container with two columns for the logo and send image
with st.container():
    col1, col2 = st.columns([1, 4])  # 1:4 ratio for the columns
    with col1:
        st.image("logo.png", width=100)
    with col2:
        st.image("send_image.png", width=600)

st.title("📄 PDF Tools Dashboard")
st.markdown("Welcome to the PDF Tools Dashboard. Choose a tool from the sidebar to get started.")

tool = st.sidebar.selectbox(
    "Select a tool", [
        "Merge PDF", "Split PDF", "Compress PDF", "Image to PDF", "Word to PDF", 
        "PDF to Word", "Unlock PDF", "Protect PDF", "PowerPoint to PDF", 
        "PDF to PowerPoint"
    ]
)

with st.sidebar.expander("Instructions", expanded=True):
    st.markdown("""
    - **Merge PDF**: Combine multiple PDF files into one.
    - **Split PDF**: Extract specific pages from a PDF.
    - **Compress PDF**: Reduce the file size of a PDF.
    - **Image to PDF**: Convert image files to a single PDF.
    - **Word to PDF**: Convert a Word document to PDF.
    - **PDF to Word**: Convert a PDF to a Word document.
    - **Unlock PDF**: Remove password protection from a PDF.
    - **Protect PDF**: Add password protection to a PDF.
    - **PowerPoint to PDF**: Convert a PowerPoint presentation to PDF.
    - **PDF to PowerPoint**: Convert a PDF to a PowerPoint presentation.
    """)

if tool == "Merge PDF":
    st.header("Merge PDF Files")
    uploaded_files = st.file_uploader("Upload PDF files", accept_multiple_files=True, type="pdf")
    if st.button("Merge"):
        if uploaded_files:
            with st.spinner("Merging PDFs..."):
                output = merge_pdfs(uploaded_files)
            st.success("PDFs merged successfully!")
            st.download_button("Download Merged PDF", output, file_name="merged.pdf")
        else:
            st.warning("Please upload PDF files to merge.")

elif tool == "Split PDF":
    st.header("Split PDF File")
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    start_page, end_page = st.columns(2)
    start_page = start_page.number_input("Start page", min_value=1)
    end_page = end_page.number_input("End page", min_value=1)
    if st.button("Split"):
        if uploaded_file and start_page <= end_page:
            with st.spinner("Splitting PDF..."):
                output = split_pdf(uploaded_file, start_page, end_page)
            st.success("PDF split successfully!")
            st.download_button("Download Split PDF", output, file_name="split.pdf")
        else:
            st.warning("Please upload a PDF file and specify valid start and end pages.")

elif tool == "Compress PDF":
    st.header("Compress PDF File")
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    if st.button("Compress"):
        if uploaded_file:
            with st.spinner("Compressing PDF..."):
                output = compress_pdf(uploaded_file)
            st.success("PDF compressed successfully!")
            st.download_button("Download Compressed PDF", output, file_name="compressed.pdf")
        else:
            st.warning("Please upload a PDF file.")

elif tool == "Image to PDF":
    st.header("Convert Images to PDF")
    uploaded_files = st.file_uploader("Upload image files", accept_multiple_files=True, type=["jpg", "jpeg", "png", "bmp", "gif"])
    if st.button("Convert"):
        if uploaded_files:
            with st.spinner("Converting images to PDF..."):
                output = images_to_pdf(uploaded_files)
            st.success("Images converted to PDF successfully!")
            st.download_button("Download PDF", output, file_name="converted.pdf")
        else:
            st.warning("Please upload image files to convert to PDF.")

elif tool == "Word to PDF":
    st.header("Convert Word to PDF")
    uploaded_file = st.file_uploader("Upload a Word file", type="docx")
    if st.button("Convert"):
        if uploaded_file:
            with st.spinner("Converting Word to PDF..."):
                output = word_to_pdf(uploaded_file)
            st.success("Word file converted to PDF successfully!")
            st.download_button("Download PDF", output, file_name="converted.pdf")
        else:
            st.warning("Please upload a Word file.")

elif tool == "PDF to Word":
    st.header("Convert PDF to Word")
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    if st.button("Convert"):
        if uploaded_file:
            with st.spinner("Converting PDF to Word..."):
                output = pdf_to_word(uploaded_file)
            st.success("PDF converted to Word successfully!")
            st.download_button("Download Word", output, file_name="converted.docx")
        else:
            st.warning("Please upload a PDF file.")

elif tool == "Unlock PDF":
    st.header("Unlock PDF File")
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    password = st.text_input("Enter password", type="password")
    if st.button("Unlock"):
        if uploaded_file and password:
            with st.spinner("Unlocking PDF..."):
                output = unlock_pdf(uploaded_file, password)
            st.success("PDF unlocked successfully!")
            st.download_button("Download Unlocked PDF", output, file_name="unlocked.pdf")
        else:
            st.warning("Please upload a PDF file and enter the password.")

elif tool == "Protect PDF":
    st.header("Protect PDF File")
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    password = st.text_input("Enter password", type="password")
    if st.button("Protect"):
        if uploaded_file and password:
            with st.spinner("Protecting PDF..."):
                output = protect_pdf(uploaded_file, password)
            st.success("PDF protected successfully!")
            st.download_button("Download Protected PDF", output, file_name="protected.pdf")
        else:
            st.warning("Please upload a PDF file and enter a password.")

elif tool == "PowerPoint to PDF":
    st.header("Convert PowerPoint to PDF")
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type="pptx")
    if st.button("Convert"):
        if uploaded_file:
            with st.spinner("Converting PowerPoint to PDF..."):
                output = ppt_to_pdf(uploaded_file)
            st.success("PowerPoint file converted to PDF successfully!")
            st.download_button("Download PDF", output, file_name="converted.pdf")
        else:
            st.warning("Please upload a PowerPoint file.")

elif tool == "PDF to PowerPoint":
    st.header("Convert PDF to PowerPoint")
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    if st.button("Convert"):
        if uploaded_file:
            with st.spinner("Converting PDF to PowerPoint..."):
                output = pdf_to_ppt(uploaded_file)
            st.success("PDF converted to PowerPoint successfully!")
            st.download_button("Download PowerPoint", output, file_name="converted.pptx")
        else:
            st.warning("Please upload a PDF file.")

# Add footer
st.markdown("""
    <style>
        .footer {
            position: fixed;
            left: 0;
            bottom: 0;
            width: 100%;
            background-color: white;
            color: black;
            text-align: center;
            padding: 10px;
        }
    </style>
    <div class="footer">
        <p>Powered by OneAmit</p>
    </div>
    """, unsafe_allow_html=True)
